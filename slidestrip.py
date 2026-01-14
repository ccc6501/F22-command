"""
slidestrip.py (REPLACEMENT)
Slide Pipeline Tool

What it does:
1) Exports every PPT slide to an image (Windows: COM automation via PowerPoint)
2) Extracts every table on every slide (cross-platform: python-pptx)
   - Optional: export images from a PDF version of the deck (no PowerPoint needed)
3) Writes:
   - slides/images/slide_0001.png ...
   - slides/manifest_slides.json   (slide_index -> image + parsed zone/sequence + title)
   - data/master_parts.json        (flattened records from all slide tables)
   - data/master_parts.csv
   - data/master_parts.xlsx

Usage:
  python slidestrip.py --pptx "Presentation7.pptx" --out "C:\\path\\project_root"
  python slidestrip.py --pptx "Presentation7.pptx" --pdf "Presentation7.pptx.pdf" --out "C:\\path\\project_root"
  python slidestrip.py --pptx "Presentation7.pptx" --out "C:\\path\\project_root" --append
  python slidestrip.py --pptx "Presentation7.pptx" --pdf "Presentation7.pptx.pdf" --out "C:\\path\\project_root" --crop-main
  python slidestrip.py --pptx "Presentation7.pptx" --pdf "Presentation7.pptx.pdf" --out "C:\\path\\project_root" --crop-main --crop-min-y 0.35
  python slidestrip.py --gui

Notes:
- Image export requires Windows PowerPoint installed (COM).
- PDF image export requires pymupdf (fitz).
- Table extraction works anywhere python-pptx works.
- XLSX output requires xlsxwriter.
- Cropped image export requires Pillow (installed with python-pptx).
"""

import os
import re
import json
import csv
import hashlib
import argparse
import sys
from datetime import datetime

# Table extraction
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

# Optional Windows image export
try:
    import comtypes.client  # type: ignore
    HAS_COMTYPES = True
except Exception:
    HAS_COMTYPES = False

# Optional XLSX export
try:
    import xlsxwriter  # type: ignore
    HAS_XLSXWRITER = True
except Exception:
    HAS_XLSXWRITER = False

# Optional PDF image export
try:
    import fitz  # type: ignore
    HAS_PYMUPDF = True
except Exception:
    HAS_PYMUPDF = False

# Optional image cropping
try:
    from PIL import Image  # type: ignore
    HAS_PIL = True
except Exception:
    HAS_PIL = False


ZONE_SEQ_PATTERNS = [
    # Common patterns you can extend:
    # "Zone 1 - Detail 3", "Z1 D3", "Z01-03", "Zone 01 Seq 03"
    re.compile(r"\bzone\s*(\d+)\b.*?\b(?:seq|sequence|detail|view)\s*(\d+)\b", re.IGNORECASE),
    re.compile(r"\bz\s*0*(\d+)\b.*?\b(?:d|s|seq|detail)\s*0*(\d+)\b", re.IGNORECASE),
    re.compile(r"\bZ\s*0*(\d+)\s*[-_]\s*0*(\d+)\b"),
]

def safe_mkdir(path: str) -> None:
    os.makedirs(path, exist_ok=True)

def sha1_text(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()

def sanitize_source_tag(path: str) -> str:
    base = os.path.splitext(os.path.basename(path))[0]
    tag = re.sub(r"[^A-Za-z0-9]+", "_", base).strip("_")
    return tag or "SRC"

def ensure_unique_slide_id(slide_id: str, existing_ids: set, source_tag: str | None) -> str:
    if slide_id not in existing_ids:
        existing_ids.add(slide_id)
        return slide_id

    tag = source_tag or "SRC"
    candidate = f"{tag}_{slide_id}"
    if candidate not in existing_ids:
        existing_ids.add(candidate)
        return candidate

    counter = 2
    while True:
        candidate = f"{tag}_{slide_id}_{counter}"
        if candidate not in existing_ids:
            existing_ids.add(candidate)
            return candidate
        counter += 1

def max_slide_index_in_images(images_dir: str) -> int:
    max_idx = 0
    if not os.path.isdir(images_dir):
        return max_idx
    for name in os.listdir(images_dir):
        m = re.match(r"slide_(\d+)\.png$", name, re.IGNORECASE)
        if m:
            max_idx = max(max_idx, int(m.group(1)))
    return max_idx

def is_picture_shape(shape) -> bool:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    except Exception:
        return False
    return False

def is_title_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        ptype = shape.placeholder_format.type
        return ptype in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE)
    except Exception:
        return False

def shape_bbox(shape):
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except Exception:
        return None
    if width <= 0 or height <= 0:
        return None
    return (left, top, left + width, top + height)

def union_bbox(a, b):
    return (min(a[0], b[0]), min(a[1], b[1]), max(a[2], b[2]), max(a[3], b[3]))

def intersects(a, b) -> bool:
    return not (a[2] < b[0] or a[0] > b[2] or a[3] < b[1] or a[1] > b[3])

def expand_bbox(b, pad_x, pad_y):
    return (b[0] - pad_x, b[1] - pad_y, b[2] + pad_x, b[3] + pad_y)

def clamp_bbox(b, max_w, max_h):
    return (max(0, b[0]), max(0, b[1]), min(max_w, b[2]), min(max_h, b[3]))

def is_header_element(shape, slide_w, slide_h) -> bool:
    """
    Detect header elements - anything in the top portion of the slide.
    This includes: Kit # title, small F-22 thumbnail, any top decorations.
    """
    bbox = shape_bbox(shape)
    if not bbox:
        return False
    
    # Get the bottom edge of this shape
    shape_bottom = bbox[3]
    
    # If the shape's bottom edge is in the top 15% of the slide, it's header
    if shape_bottom <= slide_h * 0.15:
        return True
    
    # Also check for Kit # pattern specifically (can be slightly lower)
    if getattr(shape, "has_text_frame", False):
        text = (getattr(shape, "text", "") or "").strip().upper()
        if re.search(r"KIT\s*#", text):
            return True
        # Title placeholder
        if is_title_placeholder(shape):
            return True
    
    return False


def is_footer_text(shape, slide_w, slide_h) -> bool:
    """
    Detect footer text shapes - disclaimers, reference text at bottom of slide.
    Looks for: REFERENCE, PROPRIETARY, INFORMATION, etc.
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    bbox = shape_bbox(shape)
    if not bbox:
        return False
    
    # Only consider shapes in the bottom 15% of the slide
    shape_top = bbox[1]
    if shape_top < slide_h * 0.85:
        return False
    
    text = (getattr(shape, "text", "") or "").strip().upper()
    
    # Check for common footer patterns
    footer_keywords = ["REFERENCE", "PROPRIETARY", "INFORMATION", "LOCKHEED", "MARTIN", "CONFIDENTIAL"]
    for keyword in footer_keywords:
        if keyword in text:
            return True
    
    return False


def compute_main_bbox(
    slide,
    slide_w,
    slide_h,
    pad_ratio: float = 0.01,
    min_y_ratio: float = 0.35,
    max_area_ratio: float = 0.85,
):
    """
    Compute the main content bounding box by finding the region between
    the header (Kit # title, thumbnail) and footer (REFERENCE disclaimer) areas.
    
    Strategy:
    1. Find the bottom edge of all header elements (top of slide)
    2. Find the top edge of all footer elements (bottom of slide)
    3. Crop the full width between these two boundaries
    """
    shapes = list(slide.shapes)
    
    # Find the bottom of the header area (where header elements end)
    header_bottom = 0
    for shape in shapes:
        if is_header_element(shape, slide_w, slide_h) or getattr(shape, "has_table", False):
            bbox = shape_bbox(shape)
            if bbox:
                header_bottom = max(header_bottom, bbox[3])
    
    # Find the top of the footer area (where footer text begins)
    footer_top = slide_h
    for shape in shapes:
        if is_footer_text(shape, slide_w, slide_h):
            bbox = shape_bbox(shape)
            if bbox:
                footer_top = min(footer_top, bbox[1])
    
    # If no header found, use a default top margin (12% of slide)
    if header_bottom == 0:
        header_bottom = int(slide_h * 0.12)
    
    # If no footer found, use a default bottom margin (100% - no crop)
    if footer_top >= slide_h:
        footer_top = slide_h
    
    # Ensure we have a valid region (header_bottom < footer_top)
    if header_bottom >= footer_top:
        # Fallback: use middle portion of slide
        header_bottom = int(slide_h * 0.12)
        footer_top = int(slide_h * 0.95)
    
    # Add small padding
    pad_y = int(slide_h * pad_ratio)
    pad_x = int(slide_w * pad_ratio)
    
    # Create bbox: full width, from below header to above footer
    crop_top = header_bottom + pad_y
    crop_bottom = footer_top - pad_y
    crop_left = pad_x
    crop_right = slide_w - pad_x
    
    # Clamp to valid bounds
    crop_top = max(0, crop_top)
    crop_bottom = min(slide_h, crop_bottom)
    crop_left = max(0, crop_left)
    crop_right = min(slide_w, crop_right)
    
    # Ensure minimum size
    if crop_bottom <= crop_top or crop_right <= crop_left:
        return None
    
    return (crop_left, crop_top, crop_right, crop_bottom)


# =============================================================================
# PPTX Shape-Based Export (extracts images + labels directly from shapes)
# =============================================================================

def is_part_label(shape, slide_w, slide_h) -> bool:
    """
    Detect part number labels - small text boxes with part IDs like 5HY03205-111A.
    These are the labels we want to keep in the composite.
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    
    text = (getattr(shape, "text", "") or "").strip()
    if not text:
        return False
    
    # Part numbers typically match patterns like: 5HY03205-111A, 5HY03210-104B
    # General pattern: alphanumeric with dashes, ending in letter suffix
    if re.match(r"^\d*[A-Z]{1,3}\d+-\d+[A-Z]?$", text.upper()):
        return True
    
    # Also match if it looks like a part ID (short, has numbers and dashes)
    if len(text) < 20 and re.search(r"\d+-\d+", text):
        return True
    
    return False


def is_main_diagram(shape, slide_w, slide_h) -> bool:
    """
    Detect the main diagram/blueprint image.
    This is typically the largest picture shape in the content area.
    """
    if not is_picture_shape(shape):
        return False
    
    bbox = shape_bbox(shape)
    if not bbox:
        return False
    
    # Must be reasonably large (at least 20% of slide area)
    shape_area = (bbox[2] - bbox[0]) * (bbox[3] - bbox[1])
    slide_area = slide_w * slide_h
    if shape_area < slide_area * 0.15:
        return False
    
    # Should not be in the header zone (top 12%)
    center_y = (bbox[1] + bbox[3]) / 2
    if center_y < slide_h * 0.12:
        return False
    
    return True


def is_connector_line(shape) -> bool:
    """
    Detect connector/callout lines pointing from labels to parts.
    """
    try:
        # Connectors and lines
        if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
            return True
        # Check for connector type
        shape_type_val = int(shape.shape_type)
        # 9 = msoShapeTypeLine, 21 = msoConnector
        if shape_type_val in (9, 21):
            return True
    except Exception:
        pass
    return False


def is_excluded_shape(shape, slide_w, slide_h) -> bool:
    """
    Determine if a shape should be excluded from the composite.
    Excludes: tables, Kit# title, mini aircraft, reference footer.
    """
    # Tables - always exclude
    if getattr(shape, "has_table", False):
        return True
    
    # Header elements (Kit #, mini aircraft thumbnail)
    if is_header_element(shape, slide_w, slide_h):
        return True
    
    # Footer elements (REFERENCE, PROPRIETARY, etc.)
    if is_footer_text(shape, slide_w, slide_h):
        return True
    
    return False


def extract_shape_image(shape):
    """
    Extract the image blob from a picture shape.
    Returns (image_bytes, content_type) or (None, None) if not available.
    """
    try:
        image = shape.image
        return image.blob, image.content_type
    except Exception:
        return None, None


def emu_to_pixels(emu, dpi=96):
    """Convert EMU (English Metric Units) to pixels at given DPI."""
    # 1 inch = 914400 EMU, so pixels = emu / 914400 * dpi
    return int(emu * dpi / 914400)


def export_composite_from_pptx(
    prs,
    images_dir: str,
    start_index: int = 1,
    dpi: int = 150,
    background_color: tuple = (255, 255, 255),
):
    """
    Export slide images by compositing shapes directly from PPTX.
    Includes: main diagram + part labels + connector lines
    Excludes: tables, Kit# title, mini aircraft, reference footer
    
    Args:
        prs: python-pptx Presentation object
        images_dir: Output directory for images
        start_index: Starting slide number
        dpi: Resolution for output images
        background_color: RGB tuple for background
    """
    if not HAS_PIL:
        raise RuntimeError("Pillow not available. Install Pillow for shape compositing.")
    
    from PIL import ImageDraw, ImageFont
    from io import BytesIO
    
    safe_mkdir(images_dir)
    
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    
    # Convert slide dimensions to pixels
    img_w = emu_to_pixels(slide_w_emu, dpi)
    img_h = emu_to_pixels(slide_h_emu, dpi)
    
    labels_data = []  # Collect label info for JSON output
    
    for i, slide in enumerate(prs.slides):
        slide_index = (start_index - 1) + (i + 1)
        
        # Create blank canvas
        canvas = Image.new("RGB", (img_w, img_h), background_color)
        
        slide_labels = []
        shapes_to_draw = []
        
        # First pass: categorize shapes
        for shape in slide.shapes:
            if is_excluded_shape(shape, slide_w_emu, slide_h_emu):
                continue
            
            bbox = shape_bbox(shape)
            if not bbox:
                continue
            
            # Convert EMU bbox to pixel coordinates
            px_left = emu_to_pixels(bbox[0], dpi)
            px_top = emu_to_pixels(bbox[1], dpi)
            px_right = emu_to_pixels(bbox[2], dpi)
            px_bottom = emu_to_pixels(bbox[3], dpi)
            px_bbox = (px_left, px_top, px_right, px_bottom)
            
            if is_main_diagram(shape, slide_w_emu, slide_h_emu):
                shapes_to_draw.append(("diagram", shape, px_bbox))
            elif is_part_label(shape, slide_w_emu, slide_h_emu):
                text = (shape.text or "").strip()
                shapes_to_draw.append(("label", shape, px_bbox, text))
                slide_labels.append({
                    "text": text,
                    "bbox": [px_left, px_top, px_right, px_bottom],
                    "center": [(px_left + px_right) // 2, (px_top + px_bottom) // 2]
                })
            elif is_connector_line(shape):
                shapes_to_draw.append(("line", shape, px_bbox))
            elif is_picture_shape(shape):
                # Other pictures (might be part of diagram)
                shapes_to_draw.append(("picture", shape, px_bbox))
        
        # Second pass: draw shapes onto canvas
        for item in shapes_to_draw:
            shape_type = item[0]
            shape = item[1]
            px_bbox = item[2]
            
            if shape_type in ("diagram", "picture"):
                img_bytes, content_type = extract_shape_image(shape)
                if img_bytes:
                    try:
                        shape_img = Image.open(BytesIO(img_bytes))
                        # Resize to fit bbox
                        target_w = px_bbox[2] - px_bbox[0]
                        target_h = px_bbox[3] - px_bbox[1]
                        if target_w > 0 and target_h > 0:
                            shape_img = shape_img.resize((target_w, target_h), Image.LANCZOS)
                            # Handle transparency
                            if shape_img.mode == "RGBA":
                                canvas.paste(shape_img, (px_bbox[0], px_bbox[1]), shape_img)
                            else:
                                canvas.paste(shape_img, (px_bbox[0], px_bbox[1]))
                    except Exception as e:
                        print(f"  Warning: Could not paste image shape: {e}")
            
            elif shape_type == "label":
                text = item[3]
                draw = ImageDraw.Draw(canvas)
                # Draw label box background
                draw.rectangle(px_bbox, fill=(255, 255, 255), outline=(100, 100, 100), width=1)
                # Draw text
                try:
                    # Try to use a reasonable font size based on box height
                    box_h = px_bbox[3] - px_bbox[1]
                    font_size = max(10, min(box_h - 4, 14))
                    try:
                        font = ImageFont.truetype("arial.ttf", font_size)
                    except Exception:
                        font = ImageFont.load_default()
                    
                    # Center text in box
                    text_bbox = draw.textbbox((0, 0), text, font=font)
                    text_w = text_bbox[2] - text_bbox[0]
                    text_h = text_bbox[3] - text_bbox[1]
                    text_x = px_bbox[0] + (px_bbox[2] - px_bbox[0] - text_w) // 2
                    text_y = px_bbox[1] + (px_bbox[3] - px_bbox[1] - text_h) // 2
                    draw.text((text_x, text_y), text, fill=(0, 0, 0), font=font)
                except Exception as e:
                    print(f"  Warning: Could not draw label text '{text}': {e}")
        
        # Crop to content area (exclude header/footer regions)
        content_bbox = compute_main_bbox(slide, slide_w_emu, slide_h_emu)
        if content_bbox:
            crop_left = emu_to_pixels(content_bbox[0], dpi)
            crop_top = emu_to_pixels(content_bbox[1], dpi)
            crop_right = emu_to_pixels(content_bbox[2], dpi)
            crop_bottom = emu_to_pixels(content_bbox[3], dpi)
            # Clamp to image bounds
            crop_left = max(0, min(img_w - 1, crop_left))
            crop_top = max(0, min(img_h - 1, crop_top))
            crop_right = max(crop_left + 1, min(img_w, crop_right))
            crop_bottom = max(crop_top + 1, min(img_h, crop_bottom))
            canvas = canvas.crop((crop_left, crop_top, crop_right, crop_bottom))
            
            # Adjust label coordinates to match cropped image
            for label in slide_labels:
                label["bbox"] = [
                    label["bbox"][0] - crop_left,
                    label["bbox"][1] - crop_top,
                    label["bbox"][2] - crop_left,
                    label["bbox"][3] - crop_top,
                ]
                label["center"] = [
                    label["center"][0] - crop_left,
                    label["center"][1] - crop_top,
                ]
        
        # Save image
        output_path = os.path.join(images_dir, f"slide_{slide_index:04d}.png")
        canvas.save(output_path, "PNG")
        print(f"  Exported: {output_path} ({canvas.width}x{canvas.height})")
        
        # Collect labels for this slide
        if slide_labels:
            labels_data.append({
                "slide_index": slide_index,
                "image_file": f"slide_{slide_index:04d}.png",
                "labels": slide_labels
            })
    
    # Save labels JSON
    labels_path = os.path.join(images_dir, "labels.json")
    with open(labels_path, "w", encoding="utf-8") as f:
        json.dump(labels_data, f, indent=2)
    print(f"  Saved labels: {labels_path}")
    
    return labels_data


def crop_images_to_main_area(
    prs: Presentation,
    images_dir: str,
    start_index: int = 1,
    pad_ratio: float = 0.03,
    min_y_ratio: float = 0.35,
    max_area_ratio: float = 0.85,
):
    if not HAS_PIL:
        raise RuntimeError("Pillow not available. Install Pillow to crop images.")

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    for i in range(len(prs.slides)):
        slide = prs.slides[i]
        bbox = compute_main_bbox(
            slide,
            slide_w,
            slide_h,
            pad_ratio=pad_ratio,
            min_y_ratio=min_y_ratio,
            max_area_ratio=max_area_ratio,
        )
        if not bbox:
            continue
        slide_index = (start_index - 1) + (i + 1)
        image_path = os.path.join(images_dir, f"slide_{slide_index:04d}.png")
        if not os.path.exists(image_path):
            continue
        with Image.open(image_path) as img:
            scale_x = img.width / slide_w
            scale_y = img.height / slide_h
            left = int(bbox[0] * scale_x)
            top = int(bbox[1] * scale_y)
            right = int(bbox[2] * scale_x)
            bottom = int(bbox[3] * scale_y)
            left = max(0, min(img.width - 1, left))
            top = max(0, min(img.height - 1, top))
            right = max(left + 1, min(img.width, right))
            bottom = max(top + 1, min(img.height, bottom))
            cropped = img.crop((left, top, right, bottom))
            cropped.save(image_path)

def parse_zone_sequence(text: str):
    if not text:
        return None, None
    t = " ".join(text.split())
    for pat in ZONE_SEQ_PATTERNS:
        m = pat.search(t)
        if m:
            try:
                z = int(m.group(1))
                s = int(m.group(2))
                return z, s
            except Exception:
                continue
    return None, None

def extract_slide_title(slide) -> str:
    # Best-effort: choose the longest text run on slide as a title candidate
    best = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = (shape.text or "").strip()
        if len(txt) > len(best):
            best = txt
    return best.strip()

def extract_tables_from_slide(prs: Presentation, slide_index: int):
    slide = prs.slides[slide_index]
    tables = []
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        tbl = shape.table
        rows = []
        for r in range(len(tbl.rows)):
            row = []
            for c in range(len(tbl.columns)):
                cell = tbl.cell(r, c)
                val = (cell.text or "").strip()
                row.append(val)
            rows.append(row)
        tables.append(rows)
    return tables

def flatten_tables_to_records(slide_meta: dict, tables: list):
    """
    Flatten tables into row-records.
    Heuristic:
    - First non-empty row becomes header if it looks header-ish
    - Else we store as col_0..col_n
    """
    records = []
    for t_index, rows in enumerate(tables):
        cleaned = [r for r in rows if any((c or "").strip() for c in r)]
        if not cleaned:
            continue

        header = cleaned[0]
        body = cleaned[1:] if len(cleaned) > 1 else []

        # Detect "header-ish": at least 2 non-empty cells and mostly non-numeric
        nonempty = [c for c in header if c.strip()]
        looks_header = False
        if len(nonempty) >= 2:
            numericish = 0
            for c in nonempty:
                if re.fullmatch(r"[\d\W]+", c.strip()):
                    numericish += 1
            looks_header = numericish < max(1, len(nonempty) // 2)

        if looks_header and body:
            body = fill_down_rows(body)
            keys = [k.strip() or f"col_{i}" for i, k in enumerate(header)]
            for r in body:
                rec = dict(slide_meta)
                rec["table_index"] = t_index
                for i, k in enumerate(keys):
                    rec[k] = (r[i].strip() if i < len(r) else "")
                records.append(rec)
        else:
            cleaned = fill_down_rows(cleaned)
            # No header detected -- store with col_*
            for r in cleaned:
                rec = dict(slide_meta)
                rec["table_index"] = t_index
                for i, val in enumerate(r):
                    rec[f"col_{i}"] = (val or "").strip()
                records.append(rec)

    return records

def fill_down_rows(rows: list) -> list:
    """
    Carry non-empty values downward within each column.
    This helps when PPT tables leave repeated values blank.
    """
    if not rows:
        return rows

    max_cols = max(len(r) for r in rows)
    last_vals = ["" for _ in range(max_cols)]
    for r in rows:
        # Ensure row is long enough to index safely
        if len(r) < max_cols:
            r.extend([""] * (max_cols - len(r)))
        for i, val in enumerate(r):
            cell = (val or "").strip()
            if cell:
                last_vals[i] = val
            elif last_vals[i]:
                r[i] = last_vals[i]
    return rows
def export_slide_images_windows(pptx_path: str, out_dir: str, start_index: int = 1) -> None:
    """
    Exports slide images via COM PowerPoint automation.
    Saves to out_dir with name slide_0001.png etc.
    """
    if not HAS_COMTYPES:
        raise RuntimeError("comtypes not available. Install comtypes and run on Windows with PowerPoint installed.")

    safe_mkdir(out_dir)

    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    except OSError as exc:
        raise RuntimeError(
            "PowerPoint COM automation failed. Ensure Microsoft PowerPoint is installed, "
            "or run with --no-images."
        ) from exc
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)

    # 18 = ppSaveAsPNG
    # But Export works better for a folder.
    # We'll export PNGs, then rename deterministically.
    tmp_export = os.path.join(out_dir, "_tmp_export")
    safe_mkdir(tmp_export)

    presentation.Export(tmp_export, "PNG")

    presentation.Close()
    powerpoint.Quit()

    # PowerPoint exports as Slide1.PNG, Slide2.PNG, etc.
    # Normalize into slide_0001.png
    for fname in os.listdir(tmp_export):
        m = re.match(r"Slide(\d+)\.PNG$", fname, re.IGNORECASE)
        if not m:
            continue
        idx = int(m.group(1))
        src = os.path.join(tmp_export, fname)
        dst_index = (start_index - 1) + idx
        dst = os.path.join(out_dir, f"slide_{dst_index:04d}.png")
        os.replace(src, dst)

    # Cleanup folder if empty
    try:
        os.rmdir(tmp_export)
    except Exception:
        pass

def export_slide_images_pdf(pdf_path: str, out_dir: str, zoom: float = 2.0, start_index: int = 1) -> None:
    """
    Exports slide images from a PDF via PyMuPDF.
    Saves to out_dir with name slide_0001.png etc.
    """
    if not HAS_PYMUPDF:
        raise RuntimeError("pymupdf not available. Install pymupdf to export images from PDF.")

    safe_mkdir(out_dir)

    doc = fitz.open(pdf_path)
    matrix = fitz.Matrix(zoom, zoom)
    for i in range(doc.page_count):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        out_index = (start_index - 1) + (i + 1)
        out_path = os.path.join(out_dir, f"slide_{out_index:04d}.png")
        pix.save(out_path)
    doc.close()

def write_master_parts_xlsx(xlsx_path: str, records: list, columns: list) -> None:
    """
    Write records to an XLSX file with a header row and table formatting.
    """
    if not HAS_XLSXWRITER:
        raise RuntimeError("xlsxwriter not available. Install xlsxwriter to write XLSX output.")

    wb = xlsxwriter.Workbook(xlsx_path)
    ws = wb.add_worksheet("master_parts")

    for row_idx, rec in enumerate(records, start=1):
        for col_idx, name in enumerate(columns):
            ws.write(row_idx, col_idx, rec.get(name, ""))

    if columns:
        last_row = max(1, len(records))
        last_col = len(columns) - 1
        ws.add_table(
            0,
            0,
            last_row,
            last_col,
            {"style": "Table Style Light 9", "columns": [{"header": name} for name in columns]},
        )
        ws.freeze_panes(1, 0)

    wb.close()

def run_pipeline(
    pptx_path: str,
    project_root: str,
    export_images: bool = True,
    pdf_path: str | None = None,
    append: bool = False,
    source_tag: str | None = None,
    crop_main: bool = False,
    crop_pad: float = 0.03,
    crop_min_y: float = 0.35,
    crop_max_area: float = 0.85,
    composite: bool = False,
    composite_dpi: int = 150,
):
    ts = datetime.now().strftime("%Y-%m-%d_%H%M%S")

    slides_dir = os.path.join(project_root, "slides")
    images_dir = os.path.join(slides_dir, "images")
    data_dir = os.path.join(project_root, "data")

    safe_mkdir(slides_dir)
    safe_mkdir(images_dir)
    safe_mkdir(data_dir)

    manifest_path = os.path.join(slides_dir, "manifest_slides.json")
    master_json_path = os.path.join(data_dir, "master_parts.json")

    existing_manifest = None
    existing_slides = []
    existing_records = []
    existing_slide_ids = set()
    existing_max_index = 0

    if append:
        if os.path.exists(manifest_path):
            with open(manifest_path, "r", encoding="utf-8") as f:
                existing_manifest = json.load(f)
            slides = existing_manifest.get("slides", [])
            if isinstance(slides, list):
                existing_slides = list(slides)
                for s in slides:
                    sid = s.get("slide_id")
                    if sid:
                        existing_slide_ids.add(sid)
                    try:
                        existing_max_index = max(existing_max_index, int(s.get("slide_index", 0)))
                    except Exception:
                        pass
        if os.path.exists(master_json_path):
            with open(master_json_path, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    existing_records = data

    existing_max_index = max(existing_max_index, max_slide_index_in_images(images_dir))
    start_index = existing_max_index + 1 if append else 1
    if append and not source_tag:
        source_tag = sanitize_source_tag(pptx_path)

    # 1) Export slide images if possible/desired
    if export_images:
        if pdf_path:
            export_slide_images_pdf(pdf_path, images_dir, start_index=start_index)
        elif os.name == "nt" and HAS_COMTYPES:
            export_slide_images_windows(pptx_path, images_dir, start_index=start_index)
        else:
            print("WARNING: Image export skipped.")
            print("         Provide --pdf or use Windows + PowerPoint + comtypes.")

    # 2) Extract tables + build manifest
    prs = Presentation(pptx_path)
    if crop_main:
        crop_images_to_main_area(
            prs,
            images_dir,
            start_index=start_index,
            pad_ratio=crop_pad,
            min_y_ratio=crop_min_y,
            max_area_ratio=crop_max_area,
        )

    if append and existing_manifest:
        manifest = existing_manifest
        manifest["generated_at"] = ts
        manifest["slides"] = []
        sources = manifest.get("sources")
        if not isinstance(sources, list):
            sources = []
            prev_pptx = existing_manifest.get("pptx")
            if prev_pptx:
                sources.append({"pptx": prev_pptx, "generated_at": existing_manifest.get("generated_at")})
        sources.append(
            {
                "pptx": os.path.abspath(pptx_path),
                "pdf": os.path.abspath(pdf_path) if pdf_path else None,
                "generated_at": ts,
            }
        )
        manifest["sources"] = sources
    else:
        manifest = {
            "pptx": os.path.abspath(pptx_path),
            "generated_at": ts,
            "slide_count": len(prs.slides),
            "slides": []
        }

    master_records = []

    for i in range(len(prs.slides)):
        slide = prs.slides[i]
        source_slide_index = i + 1
        slide_index = (start_index - 1) + source_slide_index

        title = extract_slide_title(slide)
        zone, seq = parse_zone_sequence(title)

        image_file = os.path.join("slides", "images", f"slide_{slide_index:04d}.png")
        slide_id = f"SLIDE_{slide_index:04d}"
        if zone is not None and seq is not None:
            slide_id = f"Z{zone:02d}_S{seq:02d}"
        slide_id = ensure_unique_slide_id(slide_id, existing_slide_ids, source_tag)

        slide_meta = {
            "slide_index": slide_index,
            "slide_id": slide_id,
            "title": title,
            "zone": zone,
            "sequence": seq,
            "image_file": image_file,
            "source_pptx": os.path.abspath(pptx_path),
            "source_slide_index": source_slide_index,
        }
        if pdf_path:
            slide_meta["source_pdf"] = os.path.abspath(pdf_path)
        if source_tag:
            slide_meta["source_tag"] = source_tag

        tables = extract_tables_from_slide(prs, i)
        # Hash tables for traceability
        tables_text = json.dumps(tables, ensure_ascii=False)
        slide_meta["tables_sha1"] = sha1_text(tables_text)
        slide_meta["table_count"] = len(tables)

        # Add to manifest
        manifest["slides"].append(slide_meta)

        # Flatten to records
        records = flatten_tables_to_records(slide_meta, tables)
        master_records.extend(records)

    if append and existing_slides:
        manifest["slides"] = existing_slides + manifest["slides"]
    manifest["slide_count"] = len(manifest["slides"])

    combined_records = existing_records + master_records if append else master_records

    # Write outputs
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    with open(master_json_path, "w", encoding="utf-8") as f:
        json.dump(combined_records, f, indent=2, ensure_ascii=False)

    # CSV: union columns across records
    all_cols = set()
    for r in combined_records:
        all_cols.update(r.keys())
    all_cols = sorted(all_cols)

    master_csv_path = os.path.join(data_dir, "master_parts.csv")
    with open(master_csv_path, "w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(f, fieldnames=all_cols)
        w.writeheader()
        for r in combined_records:
            w.writerow(r)

    master_xlsx_path = os.path.join(data_dir, "master_parts.xlsx")
    if HAS_XLSXWRITER:
        write_master_parts_xlsx(master_xlsx_path, combined_records, all_cols)
    else:
        print("WARNING: XLSX output skipped: xlsxwriter not available.")

    print("OK: Pipeline complete")
    print(" -", manifest_path)
    print(" -", master_json_path)
    print(" -", master_csv_path)
    if HAS_XLSXWRITER:
        print(" -", master_xlsx_path)
    if export_images:
        print(" -", images_dir)

def launch_gui():
    try:
        import tkinter as tk
        from tkinter import filedialog, messagebox
    except Exception as exc:
        print(f"ERROR: tkinter not available: {exc}")
        return

    class PipelineGui:
        def __init__(self, root: tk.Tk):
            self.root = root
            self.root.title("Slide Pipeline Tool")
            self.root.geometry("640x420")

            self.pptx_path = tk.StringVar()
            self.pdf_path = tk.StringVar()
            self.out_dir = tk.StringVar()
            self.append_mode = tk.StringVar(value="overwrite")
            self.skip_images = tk.BooleanVar(value=False)
            self.source_tag = tk.StringVar()
            self.crop_main = tk.BooleanVar(value=False)
            self.crop_pad = tk.DoubleVar(value=0.03)
            self.crop_min_y = tk.DoubleVar(value=0.35)

            pad = {"padx": 10, "pady": 6}

            tk.Label(root, text="PPTX File (required)").pack(anchor="w", **pad)
            self._file_picker(self.pptx_path, self.browse_pptx)

            tk.Label(root, text="PDF File (optional)").pack(anchor="w", **pad)
            self._file_picker(self.pdf_path, self.browse_pdf)

            tk.Label(root, text="Output Folder (required)").pack(anchor="w", **pad)
            self._dir_picker(self.out_dir, self.browse_out)

            tk.Label(root, text="Mode").pack(anchor="w", **pad)
            mode_frame = tk.Frame(root)
            mode_frame.pack(anchor="w", **pad)
            tk.Radiobutton(mode_frame, text="Overwrite", variable=self.append_mode, value="overwrite").pack(side="left")
            tk.Radiobutton(mode_frame, text="Append", variable=self.append_mode, value="append").pack(side="left")

            tk.Label(root, text="Source Tag (optional, used for append collisions)").pack(anchor="w", **pad)
            tk.Entry(root, textvariable=self.source_tag, width=40).pack(anchor="w", padx=10)

            tk.Checkbutton(root, text="Skip image export", variable=self.skip_images).pack(anchor="w", **pad)
            tk.Checkbutton(root, text="Crop images to main picture + labels", variable=self.crop_main).pack(anchor="w", **pad)

            pad_frame = tk.Frame(root)
            pad_frame.pack(anchor="w", padx=10, pady=2)
            tk.Label(pad_frame, text="Crop padding ratio").pack(side="left")
            tk.Entry(pad_frame, textvariable=self.crop_pad, width=6).pack(side="left", padx=6)

            miny_frame = tk.Frame(root)
            miny_frame.pack(anchor="w", padx=10, pady=2)
            tk.Label(miny_frame, text="Crop min Y ratio").pack(side="left")
            tk.Entry(miny_frame, textvariable=self.crop_min_y, width=6).pack(side="left", padx=6)

            tk.Button(root, text="Run Pipeline", command=self.run).pack(fill="x", padx=10, pady=12)

        def _file_picker(self, var: tk.StringVar, browse_cb):
            frame = tk.Frame(self.root)
            frame.pack(fill="x", padx=10, pady=2)
            tk.Entry(frame, textvariable=var, width=60).pack(side="left", fill="x", expand=True)
            tk.Button(frame, text="Browse...", command=browse_cb).pack(side="right", padx=4)

        def _dir_picker(self, var: tk.StringVar, browse_cb):
            frame = tk.Frame(self.root)
            frame.pack(fill="x", padx=10, pady=2)
            tk.Entry(frame, textvariable=var, width=60).pack(side="left", fill="x", expand=True)
            tk.Button(frame, text="Browse...", command=browse_cb).pack(side="right", padx=4)

        def browse_pptx(self):
            path = filedialog.askopenfilename(
                title="Select PPTX File",
                filetypes=[("PowerPoint Files", "*.pptx")],
            )
            if path:
                self.pptx_path.set(os.path.normpath(path))

        def browse_pdf(self):
            path = filedialog.askopenfilename(
                title="Select PDF File",
                filetypes=[("PDF Files", "*.pdf")],
            )
            if path:
                self.pdf_path.set(os.path.normpath(path))

        def browse_out(self):
            path = filedialog.askdirectory(title="Select Output Folder")
            if path:
                self.out_dir.set(os.path.normpath(path))

        def run(self):
            pptx = self.pptx_path.get().strip()
            if not pptx or not os.path.exists(pptx):
                messagebox.showerror("Missing PPTX", "Please select a valid PPTX file.")
                return
            out_dir = self.out_dir.get().strip()
            if not out_dir:
                messagebox.showerror("Missing Output", "Please select an output folder.")
                return

            pdf = self.pdf_path.get().strip() or None
            append = self.append_mode.get() == "append"
            source_tag = self.source_tag.get().strip() or None
            export_images = not self.skip_images.get()

            try:
                run_pipeline(
                    pptx,
                    out_dir,
                    export_images=export_images,
                    pdf_path=pdf,
                    append=append,
                    source_tag=source_tag,
                    crop_main=self.crop_main.get(),
                    crop_pad=self.crop_pad.get(),
                    crop_min_y=self.crop_min_y.get(),
                )
            except Exception as exc:
                messagebox.showerror("Pipeline Error", str(exc))
                return

            messagebox.showinfo("Done", "Pipeline complete.")

    root = tk.Tk()
    PipelineGui(root)
    root.mainloop()

def main():
    if "--gui" in sys.argv or len(sys.argv) == 1:
        launch_gui()
        return

    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True, help="Path to PPTX")
    ap.add_argument("--pdf", help="Optional PDF path to export images without PowerPoint")
    ap.add_argument("--out", required=True, help="Project root folder")
    ap.add_argument("--no-images", action="store_true", help="Skip slide image export")
    ap.add_argument("--append", action="store_true", help="Append to existing outputs instead of overwriting")
    ap.add_argument("--source-tag", help="Optional prefix when appending to avoid slide_id collisions")
    ap.add_argument("--crop-main", action="store_true", help="Crop images to main picture + nearby labels")
    ap.add_argument("--crop-pad", type=float, default=0.03, help="Crop padding as a ratio of slide size")
    ap.add_argument("--crop-min-y", type=float, default=0.35, help="Min Y ratio for crop anchor")
    ap.add_argument("--crop-max-area", type=float, default=0.85, help="Ignore shapes larger than this area ratio")
    ap.add_argument("--composite", action="store_true", help="Extract shapes from PPTX and composite (diagram + labels only)")
    ap.add_argument("--composite-dpi", type=int, default=150, help="DPI for composite output images")
    args = ap.parse_args()

    run_pipeline(
        args.pptx,
        args.out,
        export_images=(not args.no_images),
        pdf_path=args.pdf,
        append=args.append,
        source_tag=args.source_tag,
        crop_main=args.crop_main,
        crop_pad=args.crop_pad,
        crop_min_y=args.crop_min_y,
        crop_max_area=args.crop_max_area,
        composite=args.composite,
        composite_dpi=args.composite_dpi,
    )

if __name__ == "__main__":
    main()
